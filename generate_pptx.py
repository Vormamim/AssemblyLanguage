"""
Generates EduCPU8-Teacher-Presenter.pptx — a 16:9 teacher-facing deck in the
same house style as the other Vormamim "Teacher Presenter" decks (Hello-IPO,
Networking Junior Quest, HSC Networking, Dragonfire): dark background, top
accent bar, a two-column "idea / build" layout with two accent cards on the
right.

Layout constants (exact EMU positions, font sizes, corner radii, line
weights) are copied from VormSubs/src/courses/lofidragonfire/generate_pptx.py
so this deck is pixel-for-pixel consistent with the rest of the family —
ported rather than imported since this project has to stay standalone with
no dependency on VormSubs at runtime.

The three IPO accent colours used throughout (teal/violet/orange) are the
same ones the EduCPU-8 web app itself uses for its Input/Process/Output
panels, so the deck and the tool read as one consistent story.
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

# ---------- palette (matches style.css :root tokens) ----------
BG = '0A0E13'
PANEL = '121821'
PANEL_ALT = '1A222C'
BORDER = '26323D'
TEXT = 'E7EEF4'
TEXT_DIM = '8496A6'
TEXT_FAINT = '54636F'

NET, NET_DIM = '4FD1C5', '1E4A47'      # Input
VIOLET, VIOLET_DIM = 'B98CF2', '332048'  # Process
HOST, HOST_DIM = 'F2A65A', '4A331B'    # Output
GREEN, GREEN_BG = '8FD17A', '16241A'

FOOTER_TITLE = 'DIGITAL TECHNOLOGIES · EDUCPU-8 — VISUAL ASSEMBLY SIMULATOR'

DISPLAY_BOLD = 'Segoe UI Semibold'
DISPLAY = 'Segoe UI'
MONO = 'Consolas'

SLIDE_W = Emu(12191695)
SLIDE_H = Emu(6858000)


def rgb(hexstr):
    return RGBColor.from_string(hexstr)


def add_textbox(slide, x, y, w, h, runs, font=MONO, size=12, bold=False,
                 color=TEXT_DIM, align=PP_ALIGN.LEFT, line_spacing=1.2,
                 space_after=None):
    """runs: either a plain string, or a list of (text, bold, color) tuples for one paragraph."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    if isinstance(runs, str):
        runs = [(runs, bold, color)]
    for text, rbold, rcolor in runs:
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = rbold
        r.font.color.rgb = rgb(rcolor)
    return box


def add_bullet_list(slide, x, y, w, h, bullets, accent):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.3
        p.space_after = Pt(10)
        marker = p.add_run()
        marker.text = '•  '
        marker.font.name = MONO
        marker.font.size = Pt(13.5)
        marker.font.bold = True
        marker.font.color.rgb = rgb(accent)
        body = p.add_run()
        body.text = bullet
        body.font.name = MONO
        body.font.size = Pt(13.5)
        body.font.bold = False
        body.font.color.rgb = rgb(TEXT_DIM)
    return box


def add_rect(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_oval(slide, x, y, d, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_card(slide, x, y, w, h, fill, border):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(border)
    shp.line.width = Pt(1)
    shp.adjustments[0] = 0.06
    shp.shadow.inherit = False
    return shp


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BG)


def add_footer(slide, right_text):
    add_textbox(slide, Emu(502920), Emu(6528816), Emu(6400800), Emu(274320),
                FOOTER_TITLE, font=MONO, size=9, color=TEXT_FAINT)
    add_textbox(slide, Emu(10271455), Emu(6528816), Emu(1417320), Emu(274320),
                right_text, font=MONO, size=9, color=TEXT_FAINT, align=PP_ALIGN.RIGHT)


def add_header(slide, label, accent, title, num=None):
    """Top accent bar + oval dot + eyebrow label + big title — the header
    every content slide in the family shares."""
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Emu(76200), accent)
    add_oval(slide, Emu(502920), Emu(420624), Emu(88900), accent)
    add_textbox(slide, Emu(731520), Emu(365760), Emu(8229600), Emu(274320),
                label, font=DISPLAY_BOLD, size=11, bold=True, color=accent)
    if num is not None:
        add_textbox(slide, Emu(10728655), Emu(292608), Emu(960120), Emu(502920),
                    num, font=DISPLAY, size=26, bold=True, color=accent, align=PP_ALIGN.RIGHT)
    add_textbox(slide, Emu(502920), Emu(713232), Emu(11155680), Emu(685800),
                title, font=DISPLAY, size=27, bold=True, color=TEXT)


def lesson_slide(prs, blank, label, accent, accent_dim, title, idea_bullets,
                  build_text, tryit_text, checkpoint_text, footer_text, num=None):
    """The idea / build / try-it / checkpoint four-block layout used across
    every Vormamim Teacher Presenter deck's per-topic slides."""
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, label, accent, title, num=num)

    add_textbox(slide, Emu(502920), Emu(1600200), Emu(6400800), Emu(274320),
                'THE IDEA', font=DISPLAY_BOLD, size=11, bold=True, color=TEXT_DIM)
    add_bullet_list(slide, Emu(502920), Emu(1965960), Emu(6400800), Emu(2377440),
                     idea_bullets, accent)

    add_textbox(slide, Emu(502920), Emu(3977639), Emu(6400800), Emu(274320),
                'THE BUILD', font=DISPLAY_BOLD, size=11, bold=True, color=TEXT_DIM)
    add_textbox(slide, Emu(502920), Emu(4315968), Emu(6400800), Emu(1280160),
                build_text, font=MONO, size=12.5, color=TEXT_DIM, line_spacing=1.4)

    add_card(slide, Emu(7178040), Emu(1600200), Emu(4480560), Emu(1920240), GREEN_BG, GREEN)
    add_textbox(slide, Emu(7406640), Emu(1783080), Emu(4023360), Emu(274320),
                '\U0001f9ea TRY IT YOURSELF', font=DISPLAY_BOLD, size=10.5, bold=True, color=GREEN)
    add_textbox(slide, Emu(7406640), Emu(2103120), Emu(4023360), Emu(1325880),
                tryit_text, font=MONO, size=11.5, color=TEXT, line_spacing=1.35)

    add_card(slide, Emu(7178040), Emu(3703320), Emu(4480560), Emu(2194560), accent_dim, accent)
    add_textbox(slide, Emu(7406640), Emu(3886200), Emu(4023360), Emu(274320),
                'CHECKPOINT', font=DISPLAY_BOLD, size=10.5, bold=True, color=accent)
    add_textbox(slide, Emu(7406640), Emu(4206240), Emu(4023360), Emu(1600200),
                checkpoint_text, font=MONO, size=12.5, color=TEXT, line_spacing=1.4)

    add_footer(slide, footer_text)
    return slide


def add_table(slide, x, y, w, h, headers, rows, accent):
    gfx = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h)
    table = gfx.table

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(PANEL_ALT)
        cell.margin_left = cell.margin_right = Emu(91440)
        cell.margin_top = cell.margin_bottom = Emu(45720)
        cell.vertical_anchor = 3  # MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = header
        r.font.name = DISPLAY_BOLD
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = rgb(accent)

    for ri, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(ri + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(PANEL)
            cell.margin_left = cell.margin_right = Emu(91440)
            cell.margin_top = cell.margin_bottom = Emu(36576)
            cell.vertical_anchor = 3
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = value
            r.font.name = MONO if c == 0 else DISPLAY
            r.font.size = Pt(13)
            r.font.bold = (c == 0)
            r.font.color.rgb = rgb(NET) if c == 0 else rgb(TEXT_DIM if c else TEXT)

    return gfx


IPO_CARDS = [
    (NET, NET_DIM, 'INPUT',
     'Data goes in.',
     'A keyboard, a barcode scanner, a temperature sensor.',
     'The 16 memory cells you set before running.'),
    (VIOLET, VIOLET_DIM, 'PROCESS',
     'The system acts on the data, one instruction at a time.',
     "A calculator adding two numbers; a traffic light's timer.",
     'The program, using ACC as a scratchpad.'),
    (HOST, HOST_DIM, 'OUTPUT',
     'A result comes out.',
     'A receipt printing; a speaker playing a sound.',
     'Anything OUT or OUTC prints to the console.'),
]


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ---------- 1. Title slide ----------
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Emu(76200), NET)
    add_oval(slide, Emu(548640), Emu(777240), Emu(88900), NET)
    add_textbox(slide, Emu(777240), Emu(713232), Emu(7315200), Emu(274320),
                'DIGITAL TECHNOLOGIES · TEACHER PRESENTER', font=DISPLAY_BOLD, size=12, bold=True, color=NET)
    add_textbox(slide, Emu(548640), Emu(1371600), Emu(10972800), Emu(1463040),
                'EduCPU-8', font=DISPLAY, size=44, bold=True, color=TEXT)
    add_textbox(slide, Emu(548640), Emu(2148840), Emu(10972800), Emu(731520),
                'A Visual Model of How a Computer Thinks', font=DISPLAY, size=26, bold=True, color=NET)
    add_textbox(slide, Emu(548640), Emu(2880360), Emu(9875520), Emu(822960),
                'A made-up teaching chip that runs entirely in the browser. Students click Step and watch '
                'Input become Process become Output, one instruction at a time.',
                font=MONO, size=15, color=TEXT_DIM, line_spacing=1.3)

    add_card(slide, Emu(548640), Emu(3657600), Emu(5212080), Emu(1188720), PANEL, NET)
    add_textbox(slide, Emu(822960), Emu(3803904), Emu(4663440), Emu(274320),
                'THE BIG IDEA', font=DISPLAY_BOLD, size=13, bold=True, color=NET)
    add_textbox(slide, Emu(822960), Emu(4169664), Emu(4663440), Emu(548640),
                'Every computer, from a calculator to a phone, repeats the same Input → Process → Output cycle. '
                'This deck and tool make that cycle visible and clickable.',
                font=MONO, size=10.5, color=TEXT_DIM, line_spacing=1.3)

    add_card(slide, Emu(5989320), Emu(3657600), Emu(5623560), Emu(1188720), PANEL, HOST)
    add_textbox(slide, Emu(6263640), Emu(3803904), Emu(5120640), Emu(274320),
                'THE TOOL', font=DISPLAY_BOLD, size=13, bold=True, color=HOST)
    add_textbox(slide, Emu(6263640), Emu(4169664), Emu(5120640), Emu(548640),
                'A 9-instruction fake chip with 16 memory cells — small enough to fit on one screen, '
                'built to be watched rather than just written.',
                font=MONO, size=10.5, color=TEXT_DIM, line_spacing=1.3)

    add_card(slide, Emu(548640), Emu(5074920), Emu(11064240), Emu(1051560), PANEL_ALT, BORDER)
    add_textbox(slide, Emu(822960), Emu(5239512), Emu(10515600), Emu(274320),
                'RUNS ENTIRELY IN THE BROWSER — NO INSTALL, NO ACCOUNTS',
                font=DISPLAY_BOLD, size=10.5, bold=True, color=TEXT_FAINT)
    add_textbox(slide, Emu(822960), Emu(5532120), Emu(10515600), Emu(457200),
                'Open index.html and go. Three ready-made example programs are one click away; the '
                'Theory & Help panel (top right of the tool) carries this same reference.',
                font=MONO, size=12, color=TEXT_DIM, line_spacing=1.3)

    add_footer(slide, 'Title')

    # ---------- 2. The IPO model ----------
    lesson_slide(
        prs, blank,
        label='THE BIG IDEA', accent=NET, accent_dim=NET_DIM,
        title='Every Computer Does the Same Three Things',
        idea_bullets=[
            'Every computer — a calculator, a phone, this browser tab — repeats the same cycle: Input, Process, Output.',
            'Input is data going in; Process is the system acting on it; Output is the result coming out.',
            'The three panels on the EduCPU-8 page are colour-coded to match: teal Input, violet Process, orange Output.',
        ],
        build_text='On the EduCPU-8 page, Input is the memory grid you edit before running, Process is the '
                    'program plus the ACC / PC / Zero registers, and Output is the console that OUT and '
                    'OUTC print to.',
        tryit_text='Before showing the tool, ask students to name the Input / Process / Output steps in '
                    'something ordinary — a vending machine, a microwave, an ATM.',
        checkpoint_text="If a calculator's Input is the buttons you press, what's its Output — and where "
                         'does the Process actually happen?',
        footer_text='02',
        num='02',
    )

    # ---------- 3. Meet EduCPU-8 ----------
    lesson_slide(
        prs, blank,
        label='THE TOOL', accent=VIOLET, accent_dim=VIOLET_DIM,
        title='A Tiny Made-Up Chip, Built to Be Watched',
        idea_bullets=[
            'ACC is the only general register — every value passes through it, so "processing" is always visible in one place.',
            '16 memory cells (addresses 0–15) hold data; the program is a separate list of instructions, so there’s no self-modifying code to explain.',
            'PC always shows which line runs next; the Zero flag records whether the last LOAD/ADD/SUB produced 0.',
        ],
        build_text='Nine instructions cover everything: moving data (LOAD/STORE), arithmetic (ADD/SUB), '
                    'control flow (JMP/JZ), output (OUT/OUTC), and stopping (HALT). The full reference is '
                    'one click away in the Theory & Help panel — and on the next slide.',
        tryit_text='Open the Theory & Help slide-out with students on screen — it has the same '
                    'instruction table as the next slide, so they know where to find it again.',
        checkpoint_text='Why does giving the chip only one general register (ACC) make it easier to '
                         'follow what’s happening, instead of harder?',
        footer_text='03',
        num='03',
    )

    # ---------- 4. Instruction set reference ----------
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, 'REFERENCE', HOST, 'Instruction Set Reference', num='04')
    add_textbox(slide, Emu(502920), Emu(1554480), Emu(11155680), Emu(365760),
                'Nine instructions. Every EduCPU-8 program is built from these, and nothing else.',
                font=MONO, size=13, color=TEXT_DIM, line_spacing=1.3)
    add_table(
        slide, Emu(502920), Emu(2011680), Emu(11155680), Emu(4297680),
        headers=['OP', 'OPERAND', 'WHAT IT DOES'],
        rows=[
            ['LOAD', 'addr or #val', 'ACC ← memory[addr], or ACC ← val'],
            ['STORE', 'addr', 'memory[addr] ← ACC'],
            ['ADD', 'addr or #val', 'ACC ← ACC + (…)'],
            ['SUB', 'addr or #val', 'ACC ← ACC − (…)'],
            ['JMP', 'line', 'jump to that line'],
            ['JZ', 'line', 'jump there only if ACC is 0'],
            ['OUT', '—', 'print ACC as a number'],
            ['OUTC', '—', "print ACC as a character — #'H' is shorthand for a character's code"],
            ['HALT', '—', 'stop the program'],
        ],
        accent=HOST,
    )
    add_footer(slide, '04')

    # ---------- 5, 6, 7. Walkthroughs ----------
    lesson_slide(
        prs, blank,
        label='WALKTHROUGH · OUTPUT', accent=HOST, accent_dim=HOST_DIM,
        title='Hello Text',
        idea_bullets=[
            'No loop needed — every character is printed by its own LOAD + OUTC pair.',
            "#'H' is shorthand for a character's ASCII code; LOAD #'H' then OUTC prints “H”.",
            'This program has no memory or Input step at all — it’s Output in isolation.',
        ],
        build_text="LOAD #'H' → OUTC → LOAD #'i' → OUTC → LOAD #'!' → OUTC → HALT "
                    'prints Hi! to the console, one character at a time.',
        tryit_text='Change the characters (or add more LOAD/OUTC pairs) to print a student’s own '
                    'name instead of “Hi!”.',
        checkpoint_text="Why does OUTC need a character's numeric code, and where does #'H' actually come from?",
        footer_text='05',
        num='05',
    )

    lesson_slide(
        prs, blank,
        label='WALKTHROUGH · PROCESS', accent=VIOLET, accent_dim=VIOLET_DIM,
        title='Countdown Loop',
        idea_bullets=[
            'JZ jumps only when the Zero flag is set — the loop keeps going until the counter hits 0.',
            'SUB #1 shrinks the counter by exactly one each pass through the loop.',
            'JMP unconditionally sends execution back to the top of the loop.',
        ],
        build_text='Memory cell 0 starts as the counter. Each pass prints its value, subtracts 1, stores '
                    'it back, then jumps back — until JZ sends it straight to HALT.',
        tryit_text='Click memory cell 0 before running and change the starting count — try 1, then try '
                    'a bigger number and switch to Run instead of Step.',
        checkpoint_text='What would happen if the JZ line were deleted — would the program ever stop, and why?',
        footer_text='06',
        num='06',
    )

    lesson_slide(
        prs, blank,
        label='WALKTHROUGH · INPUT', accent=NET, accent_dim=NET_DIM,
        title='Add Two Numbers',
        idea_bullets=[
            'Memory cells 0 and 1 are the Input — editable right after the example loads, before you Step or Run.',
            'LOAD then ADD combines two values through ACC — the same "everything passes through one place" idea as the chip itself.',
            'STORE writes the result back to memory before OUT prints it — Process finishing before Output begins.',
        ],
        build_text='LOAD 0 → ADD 1 → STORE 2 → OUT → HALT adds whatever is in cells 0 and 1, '
                    'saves the sum in cell 2, and prints it.',
        tryit_text='Edit memory cells 0 and 1 to two different numbers before running, and predict the '
                    'printed sum before clicking Step.',
        checkpoint_text='Which single memory cell holds this program’s result — and how is that '
                         'different to what OUT actually prints?',
        footer_text='07',
        num='07',
    )

    # ---------- 8. Suggested lesson flow ----------
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, 'IN THE CLASSROOM', HOST, 'A Suggested 45-Minute Lesson Flow', num='08')

    FLOW = [
        ('5 MIN', NET, 'Discuss IPO with an everyday example — a vending machine, a microwave, an ATM.'),
        ('10 MIN', VIOLET, 'Demo: load Hello Text together as a class and Step through it line by line.'),
        ('15 MIN', VIOLET, 'Pairs: work through Countdown Loop and Add Two Numbers, using each slide’s checkpoint question.'),
        ('10 MIN', HOST, 'Challenge: write a fourth program from scratch — print a name, or add three numbers.'),
        ('5 MIN', NET, 'Wrap-up: what changed in ACC, memory and the console at each step?'),
    ]
    row_h = Emu(658368)
    row_gap = 786384
    row_y0 = 1600200
    for i, (time_label, accent, text) in enumerate(FLOW):
        y = Emu(row_y0 + i * row_gap)
        add_card(slide, Emu(502920), y, Emu(11155680), row_h, PANEL, BORDER)
        add_rect(slide, Emu(502920), y, Emu(54864), row_h, accent)
        add_textbox(slide, Emu(731520), Emu(y.emu + 91440), Emu(1005840), Emu(457200),
                    time_label, font=DISPLAY_BOLD, size=13, bold=True, color=accent, align=PP_ALIGN.LEFT)
        add_textbox(slide, Emu(1783080), Emu(y.emu + 91440), Emu(9601200), Emu(457200),
                    text, font=MONO, size=12.5, color=TEXT, line_spacing=1.3)

    add_footer(slide, '08')

    # ---------- 9. Ten challenges for students ----------
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, 'ON YOUR OWN', VIOLET, '10 Challenges for Students', num='09')
    add_textbox(slide, Emu(502920), Emu(1554480), Emu(11155680), Emu(365760),
                'Each one builds on a worked example from this deck — same nine instructions, new problem.',
                font=MONO, size=12.5, color=TEXT_DIM, line_spacing=1.3)

    CHALLENGES = [
        ('1', 'Print your initials', 'Warm-up', NET, "LOAD #'…', OUTC"),
        ('2', 'Triple a number', 'Warm-up', NET, 'LOAD, ADD, ADD, STORE, OUT'),
        ('3', 'Is it zero?', 'Warm-up', NET, 'LOAD, JZ, branch to YES/NO'),
        ('4', 'Count up to a target', 'Core', VIOLET, 'loop, ADD #1, JZ to stop'),
        ('5', 'Sum from 1 to N', 'Core', VIOLET, 'loop + a running total in memory'),
        ('6', 'Halve a number by counting', 'Core', VIOLET, 'repeated SUB, count how many times'),
        ('7', 'Are three numbers all equal?', 'Core', VIOLET, 'two chained comparisons'),
        ('8', 'Countdown, then shout GO!', 'Core', VIOLET, 'a loop followed by fixed text output'),
        ('9', 'Predict the overflow', 'Stretch', HOST, 'hand-calculate first, then Step to check'),
        ('10', 'Design your own', 'Stretch', HOST, 'combine instructions nobody told you to combine'),
    ]

    gfx = add_table(
        slide, Emu(502920), Emu(2011680), Emu(11155680), Emu(4350000),
        headers=['#', 'CHALLENGE', 'LEVEL', 'SKILLS'],
        rows=[[num, title, level, skills] for num, title, level, _accent, skills in CHALLENGES],
        accent=VIOLET,
    )
    for col, width in zip(gfx.table.columns, [Emu(548640), Emu(4389120), Emu(1554480), Emu(4663440)]):
        col.width = width
    for ri, (_num, _title, level, accent, _skills) in enumerate(CHALLENGES):
        cell = gfx.table.cell(ri + 1, 2)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = rgb(accent)
                r.font.bold = True

    add_footer(slide, '09')

    # ---------- 10. Wrap-up ----------
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Emu(76200), VIOLET)
    add_oval(slide, Emu(548640), Emu(777240), Emu(88900), VIOLET)
    add_textbox(slide, Emu(777240), Emu(713232), Emu(7315200), Emu(274320),
                'RECAP', font=DISPLAY_BOLD, size=12, bold=True, color=VIOLET)
    add_textbox(slide, Emu(548640), Emu(1051560), Emu(10972800), Emu(822960),
                'Input. Process. Output. Every Time.', font=DISPLAY, size=32, bold=True, color=TEXT)
    add_textbox(slide, Emu(548640), Emu(1828800), Emu(10515600), Emu(457200),
                'Once students can point to Input, Process and Output on EduCPU-8, they can point to it '
                'in any system — a phone, a website, a game.',
                font=MONO, size=13, color=TEXT_DIM, line_spacing=1.3)

    card_w = Emu(3608705)
    card_gap = 219075
    card_y = Emu(2743200)
    card_h = Emu(2560320)
    for i, (accent, accent_dim, label, definition, example, in_tool) in enumerate(IPO_CARDS):
        x = Emu(548640 + i * (card_w.emu + card_gap))
        add_card(slide, x, card_y, card_w, card_h, accent_dim, accent)
        add_textbox(slide, Emu(x.emu + 228600), Emu(card_y.emu + 182880), Emu(card_w.emu - 457200), Emu(320040),
                    label, font=DISPLAY_BOLD, size=15, bold=True, color=accent)
        add_textbox(slide, Emu(x.emu + 228600), Emu(card_y.emu + 594360), Emu(card_w.emu - 457200), Emu(548640),
                    definition, font=MONO, size=12.5, bold=True, color=TEXT, line_spacing=1.3)
        add_textbox(slide, Emu(x.emu + 228600), Emu(card_y.emu + 1188720), Emu(card_w.emu - 457200), Emu(731520),
                    [('EXAMPLE: ', True, TEXT_FAINT), (example, False, TEXT_DIM)], font=MONO, size=10.5, line_spacing=1.35)
        add_textbox(slide, Emu(x.emu + 228600), Emu(card_y.emu + 1920240), Emu(card_w.emu - 457200), Emu(548640),
                    [('IN EDUCPU-8: ', True, TEXT_FAINT), (in_tool, False, TEXT_DIM)], font=MONO, size=10.5, line_spacing=1.35)

    add_textbox(slide, Emu(548640), Emu(5578475), Emu(10972800), Emu(457200),
                'The tool: open index.html in any browser — no install, no accounts.',
                font=MONO, size=12, color=TEXT_FAINT, line_spacing=1.3)

    add_footer(slide, 'Wrap-up')

    out_path = 'EduCPU8-Teacher-Presenter.pptx'
    prs.save(out_path)
    print(f'Saved {out_path} with {len(prs.slides._sldIdLst)} slides')


if __name__ == '__main__':
    build()
